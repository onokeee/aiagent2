"""PowerPointレポートの生成。会議でそのまま映せる体裁で作る。

作りの方針:
  - 1スライド1メッセージ。上部の「キーメッセージ」に結論を1行で書き、
    図表はその根拠として下に置く。読み手は上の1行だけで用が足りる。
  - 日本語フォントを明示的に指定する。指定しないと英字フォントが当たり、
    開いた瞬間に「ちゃんとしていない資料」に見える。
  - グラフはPowerPointネイティブ（編集可）を既定にし、
    ネイティブで表現できない種類だけ画像として貼る。

1スライド = 1つの dict。kind で中身が決まる:
    title    表紙
    agenda   目次
    section  中扉
    message  文字だけ（結論・考察）
    table    表
    chart    グラフ
    kpi      数字を大きく並べる
    compare  2つ並べて比較
    closing  まとめ／次のアクション
"""
from __future__ import annotations

import io
import re
from datetime import datetime

from pptx import Presentation
from pptx.chart.data import CategoryChartData, XyChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

import config

SLIDE_W = Inches(13.333)          # 16:9
SLIDE_H = Inches(7.5)

CHART_TYPES = {
    "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "bar_stacked": XL_CHART_TYPE.COLUMN_STACKED,
    "bar_percent": XL_CHART_TYPE.COLUMN_STACKED_100,
    "hbar": XL_CHART_TYPE.BAR_CLUSTERED,
    "hbar_stacked": XL_CHART_TYPE.BAR_STACKED,
    "line": XL_CHART_TYPE.LINE_MARKERS,
    "area": XL_CHART_TYPE.AREA,
    "area_stacked": XL_CHART_TYPE.AREA_STACKED,
    "pie": XL_CHART_TYPE.PIE,
    "doughnut": XL_CHART_TYPE.DOUGHNUT,
    "scatter": XL_CHART_TYPE.XY_SCATTER,
    "radar": XL_CHART_TYPE.RADAR_MARKERS,
}
SLIDE_KINDS = ("title", "agenda", "section", "message", "table", "chart",
               "kpi", "compare", "closing")

# 配色。1枚に何色も出さない。強調は1色だけ使う。
NAVY = RGBColor(0x1F, 0x3B, 0x5C)
ACCENT = RGBColor(0x2E, 0x75, 0xB6)
HILITE = RGBColor(0xC5, 0x5A, 0x11)
INK = RGBColor(0x22, 0x26, 0x2B)
MUTED = RGBColor(0x6B, 0x72, 0x80)
LINE = RGBColor(0xD5, 0xDB, 0xE2)
BAND = RGBColor(0xF5, 0xF8, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GOOD = RGBColor(0x1E, 0x7A, 0x3C)
BAD = RGBColor(0xB0, 0x2A, 0x2A)
SERIES = ["1F4E79", "F4B183", "70AD47", "C55A11", "7F7F7F",
          "2E75B6", "A9D18E", "FFD966", "9DC3E6", "BFBFBF"]

MAX_TABLE_ROWS = 12               # これを超えると字が小さくなって読めない
MAX_CATEGORIES = 24

# 余白（本文の左右端）
MARGIN = Inches(0.62)
BODY_W = SLIDE_W - MARGIN * 2


class ReportError(Exception):
    """レポートを作れない理由（そのまま画面に出す）。"""


# =============================================================================
# 文字まわり
# =============================================================================

def _jp(run):
    """日本語フォントを当てる。

    python-pptx は latin フォントしか設定しないので、
    日本語部分に別のフォントが当たってしまう。東アジア用を直接書く。
    """
    run.font.name = config.REPORT_FONT_JA
    rPr = run.font._element          # これ自体が rPr（文字の書式）
    for tag in ("a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", config.REPORT_FONT_JA)


def _text(frame, lines, *, size=18, bold=False, color=INK, align=PP_ALIGN.LEFT,
          space_after=4, line_spacing=1.25):
    """text_frame に段落を流し込む。lines は文字列か (文字列, 上書き) の並び。"""
    frame.word_wrap = True
    items = lines if isinstance(lines, (list, tuple)) else [lines]
    first = True
    for item in items:
        opts = {}
        if isinstance(item, tuple):
            item, opts = item
        for line in str(item).split("\n"):
            p = frame.paragraphs[0] if first else frame.add_paragraph()
            first = False
            p.text = line
            p.alignment = opts.get("align", align)
            p.space_after = Pt(opts.get("space_after", space_after))
            p.line_spacing = opts.get("line_spacing", line_spacing)
            if opts.get("level"):
                p.level = opts["level"]
            for run in p.runs:
                run.font.size = Pt(opts.get("size", size))
                run.font.bold = opts.get("bold", bold)
                run.font.color.rgb = opts.get("color", color)
                _jp(run)


def _box(slide, left, top, width, height, lines, **kw):
    shape = slide.shapes.add_textbox(left, top, width, height)
    _text(shape.text_frame, lines, **kw)
    return shape


def _rect(slide, left, top, width, height, fill=None, line=None,
          shape=MSO_SHAPE.RECTANGLE):
    s = slide.shapes.add_shape(shape, left, top, width, height)
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(1)
    s.shadow.inherit = False
    return s


def _clean(v) -> str:
    return "" if v is None else str(v)


def _num(v):
    if v is None:
        return None
    if isinstance(v, bool):
        return None
    if isinstance(v, (int, float)):
        return float(v)
    s = re.sub(r"[,\s¥$%]", "", str(v))
    try:
        return float(s)
    except ValueError:
        return None


def _fmt(v) -> str:
    if isinstance(v, bool):
        return "はい" if v else "いいえ"
    if isinstance(v, float):
        if v == int(v) and abs(v) < 1e15:
            return f"{int(v):,}"
        # 12.30 ではなく 12.3 と出す（末尾の0は読み手には意味が無い）
        return f"{v:,.2f}".rstrip("0").rstrip(".")
    if isinstance(v, int):
        return f"{v:,}"
    return _clean(v)


# =============================================================================
# 共通の枠（ヘッダ・キーメッセージ・フッタ）
# =============================================================================

def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _header(slide, title: str, message: str = "") -> Emu:
    """見出しと、その下のキーメッセージ帯。戻り値は本文を始めてよい上端。"""
    _box(slide, MARGIN, Inches(0.30), BODY_W, Inches(0.55), title,
         size=25, bold=True, color=NAVY)
    _rect(slide, MARGIN, Inches(0.92), BODY_W, Emu(12700), fill=NAVY)

    if not message:
        return Inches(1.18)
    # 結論を1行で。ここだけ読めば分かるようにする。
    bar = _rect(slide, MARGIN, Inches(1.06), BODY_W, Inches(0.62), fill=BAND)
    bar.line.color.rgb = LINE
    bar.line.width = Pt(0.75)
    tf = bar.text_frame
    tf.margin_left, tf.margin_right = Inches(0.16), Inches(0.16)
    tf.margin_top = tf.margin_bottom = Inches(0.04)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    _text(tf, message, size=15, bold=True, color=NAVY, space_after=0)
    return Inches(1.86)


def _footer(slide, text: str, page: int | None = None):
    _rect(slide, MARGIN, SLIDE_H - Inches(0.52), BODY_W, Emu(9525), fill=LINE)
    if text:
        _box(slide, MARGIN, SLIDE_H - Inches(0.46), Inches(9), Inches(0.32),
             text, size=9.5, color=MUTED)
    if page:
        _box(slide, SLIDE_W - MARGIN - Inches(0.8), SLIDE_H - Inches(0.46),
             Inches(0.8), Inches(0.32), str(page), size=9.5, color=MUTED,
             align=PP_ALIGN.RIGHT)


def _notes(slide, text: str):
    if text:
        slide.notes_slide.notes_text_frame.text = str(text)


def _source(slide, top, text: str):
    """出典・条件。数字の資料には必ず要る。"""
    if text:
        _box(slide, MARGIN, top, BODY_W, Inches(0.3), f"出所: {text}",
             size=9.5, color=MUTED)


# =============================================================================
# スライドの種類ごと
# =============================================================================

def _slide_title(prs, spec):
    slide = _blank(prs)
    _rect(slide, 0, 0, SLIDE_W, Inches(3.05), fill=NAVY)
    _rect(slide, 0, Inches(3.05), SLIDE_W, Inches(0.06), fill=HILITE)
    _box(slide, Inches(0.9), Inches(1.05), SLIDE_W - Inches(1.8), Inches(1.2),
         spec.get("title", "レポート"), size=40, bold=True, color=WHITE)
    if spec.get("subtitle"):
        _box(slide, Inches(0.9), Inches(2.25), SLIDE_W - Inches(1.8), Inches(0.5),
             spec["subtitle"], size=18, color=RGBColor(0xC5, 0xD5, 0xE8))

    y = Inches(3.55)
    for line in (spec.get("lines") or [])[:4]:
        _rect(slide, Inches(0.9), y + Inches(0.10), Inches(0.09), Inches(0.22),
              fill=HILITE)
        _box(slide, Inches(1.15), y, SLIDE_W - Inches(2.2), Inches(0.42), line,
             size=16, color=INK)
        y += Inches(0.52)

    org = spec.get("org") or config.REPORT_ORG
    foot = " ／ ".join(x for x in (org, spec.get("author")) if x)
    _box(slide, Inches(0.9), SLIDE_H - Inches(1.05), Inches(8), Inches(0.34),
         spec.get("date") or datetime.now().strftime("%Y年%m月%d日"),
         size=13, color=MUTED)
    if foot:
        _box(slide, Inches(0.9), SLIDE_H - Inches(0.72), Inches(8), Inches(0.34),
             foot, size=13, color=MUTED)
    _notes(slide, spec.get("notes", ""))
    return slide


def _slide_agenda(prs, spec):
    slide = _blank(prs)
    top = _header(slide, spec.get("title") or "本日の内容", spec.get("message", ""))
    items = spec.get("items") or []
    y = top + Inches(0.18)
    step = min(Inches(0.72), (SLIDE_H - y - Inches(0.9)) / max(len(items), 1))
    for i, it in enumerate(items, 1):
        label = it.get("text") if isinstance(it, dict) else str(it)
        note = it.get("note", "") if isinstance(it, dict) else ""
        n = _rect(slide, MARGIN, y, Inches(0.44), Inches(0.44), fill=NAVY,
                  shape=MSO_SHAPE.OVAL)
        tf = n.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        _text(tf, str(i), size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
              space_after=0)
        _box(slide, MARGIN + Inches(0.62), y + Inches(0.02), BODY_W - Inches(0.8),
             Inches(0.4), label, size=17, bold=True, color=INK)
        if note:
            _box(slide, MARGIN + Inches(0.62), y + Inches(0.34),
                 BODY_W - Inches(0.8), Inches(0.3), note, size=12, color=MUTED)
        y += step
    _notes(slide, spec.get("notes", ""))
    return slide


def _slide_section(prs, spec):
    slide = _blank(prs)
    _rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=NAVY)
    _rect(slide, MARGIN, Inches(3.18), Inches(0.9), Inches(0.07), fill=HILITE)
    _box(slide, MARGIN, Inches(3.42), SLIDE_W - MARGIN * 2, Inches(0.9),
         spec.get("title", ""), size=32, bold=True, color=WHITE)
    if spec.get("subtitle"):
        _box(slide, MARGIN, Inches(4.35), SLIDE_W - MARGIN * 2, Inches(0.5),
             spec["subtitle"], size=15, color=RGBColor(0xC5, 0xD5, 0xE8))
    _notes(slide, spec.get("notes", ""))
    return slide


def _bullets(slide, left, top, width, height, items, *, size=16):
    """箇条書き。dict なら {text, level, strong} を見る。"""
    y = top
    for item in items:
        if isinstance(item, dict):
            text, level = item.get("text", ""), int(item.get("level", 0))
            strong = bool(item.get("strong"))
        else:
            text, level, strong = str(item), 0, False
        if y > top + height - Inches(0.3):
            break
        mark_x = left + Inches(0.26) * level
        if level == 0:
            _rect(slide, mark_x, y + Inches(0.13), Inches(0.10), Inches(0.10),
                  fill=HILITE if strong else ACCENT)
        else:
            _box(slide, mark_x, y - Inches(0.02), Inches(0.2), Inches(0.3), "－",
                 size=13, color=MUTED)
        _box(slide, mark_x + Inches(0.24), y - Inches(0.05),
             width - Inches(0.26) * level - Inches(0.24), Inches(0.42), text,
             size=size - 1.5 * level, bold=strong,
             color=HILITE if strong else INK)
        # 行数ぶん送る（おおよそでよい。重ならなければ十分）
        per = max(1, int(width / Inches(0.135) / max(size - 1.5 * level, 1) * 1.9))
        lines = max(1, (len(text) + per - 1) // per)
        y += Inches(0.34) * lines + Inches(0.12)
    return y


def _slide_message(prs, spec):
    slide = _blank(prs)
    top = _header(slide, spec.get("title", ""), spec.get("message", ""))
    y = top + Inches(0.12)
    if spec.get("lead"):
        _box(slide, MARGIN, y, BODY_W, Inches(0.6), spec["lead"], size=17,
             color=INK)
        y += Inches(0.72)
    items = spec.get("bullets") or []
    if items:
        y = _bullets(slide, MARGIN, y, BODY_W, SLIDE_H - y - Inches(1.2), items,
                     size=17)
    if spec.get("body"):
        _box(slide, MARGIN, y + Inches(0.1), BODY_W,
             SLIDE_H - y - Inches(1.1), spec["body"], size=14, color=INK)
    if spec.get("callout"):
        _callout(slide, MARGIN, SLIDE_H - Inches(1.55), BODY_W, spec["callout"])
    _notes(slide, spec.get("notes", ""))
    return slide


def _callout(slide, left, top, width, text, label="ポイント"):
    """強調枠。1枚に1つだけ置く。"""
    box = _rect(slide, left, top, width, Inches(0.92), fill=RGBColor(0xFD, 0xF3, 0xE7))
    box.line.color.rgb = HILITE
    box.line.width = Pt(1.25)
    _rect(slide, left, top, Inches(0.07), Inches(0.92), fill=HILITE)
    _box(slide, left + Inches(0.22), top + Inches(0.09), Inches(2), Inches(0.26),
         label, size=10.5, bold=True, color=HILITE)
    _box(slide, left + Inches(0.22), top + Inches(0.34), width - Inches(0.44),
         Inches(0.52), text, size=14, color=INK)


def _slide_kpi(prs, spec):
    slide = _blank(prs)
    top = _header(slide, spec.get("title", ""), spec.get("message", ""))
    items = (spec.get("items") or [])[:4]
    if not items:
        raise ReportError("kpi スライドには items が必要です。")
    gap = Inches(0.32)
    width = (BODY_W - gap * (len(items) - 1)) / len(items)
    card_h = Inches(2.35)
    for i, it in enumerate(items):
        left = MARGIN + (width + gap) * i
        card = _rect(slide, left, top + Inches(0.25), width, card_h, fill=BAND,
                     line=LINE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        card.adjustments[0] = 0.04
        _rect(slide, left, top + Inches(0.25), width, Inches(0.07), fill=ACCENT)
        _box(slide, left, top + Inches(0.48), width, Inches(0.34),
             _clean(it.get("label")), size=13, color=MUTED, align=PP_ALIGN.CENTER)
        _box(slide, left, top + Inches(0.86), width, Inches(0.85),
             _fmt(it.get("value")) + _clean(it.get("unit")),
             size=34, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        d = _num(it.get("delta"))
        if d is not None:
            mark = "▲" if d > 0 else ("▼" if d < 0 else "―")
            good = it.get("higher_is_better", True)
            col = MUTED if d == 0 else (GOOD if (d > 0) == bool(good) else BAD)
            _box(slide, left, top + Inches(1.72), width, Inches(0.34),
                 f"{mark} {_fmt(abs(d))}{_clean(it.get('delta_unit'))}"
                 + (f"（{it['delta_label']}）" if it.get("delta_label") else ""),
                 size=13, bold=True, color=col, align=PP_ALIGN.CENTER)
        if it.get("note"):
            _box(slide, left, top + Inches(2.06), width, Inches(0.3),
                 it["note"], size=10.5, color=MUTED, align=PP_ALIGN.CENTER)

    y = top + card_h + Inches(0.45)
    if spec.get("bullets"):
        y = _bullets(slide, MARGIN, y, BODY_W, SLIDE_H - y - Inches(1.0),
                     spec["bullets"], size=15)
    elif spec.get("comment"):
        _box(slide, MARGIN, y, BODY_W, Inches(1.0), spec["comment"], size=14)
    if spec.get("callout"):
        _callout(slide, MARGIN, SLIDE_H - Inches(1.5), BODY_W, spec["callout"])
    _source(slide, SLIDE_H - Inches(0.85), spec.get("source", ""))
    _notes(slide, spec.get("notes", ""))
    return slide


def _add_table(slide, left, top, width, height, columns, rows, *,
               font=11.5, highlight_rows=()):
    shape = slide.shapes.add_table(len(rows) + 1, len(columns), left, top,
                                   width, height)
    table = shape.table
    for j, c in enumerate(columns):
        cell = table.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        cell.margin_left = cell.margin_right = Inches(0.07)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        _text(cell.text_frame, _clean(c), size=font, bold=True, color=WHITE,
              space_after=0, line_spacing=1.0)
    for i, row in enumerate(rows, start=1):
        for j in range(len(columns)):
            v = row[j] if j < len(row) else ""
            cell = table.cell(i, j)
            cell.margin_left = cell.margin_right = Inches(0.07)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            cell.fill.fore_color.rgb = (
                RGBColor(0xFD, 0xF3, 0xE7) if (i - 1) in highlight_rows
                else (BAND if i % 2 == 0 else WHITE))
            _text(cell.text_frame, _fmt(v), size=font, space_after=0,
                  line_spacing=1.0,
                  align=PP_ALIGN.RIGHT if isinstance(v, (int, float))
                  and not isinstance(v, bool) else PP_ALIGN.LEFT,
                  bold=(i - 1) in highlight_rows)
    return table


def _slide_table(prs, spec):
    slide = _blank(prs)
    top = _header(slide, spec.get("title", ""), spec.get("message", ""))
    cols = [_clean(c) for c in (spec.get("columns") or [])]
    rows = spec.get("rows") or []
    if not cols:
        raise ReportError("table スライドには columns が必要です。")
    limit = int(spec.get("max_rows") or MAX_TABLE_ROWS)
    shown, cut = rows[:limit], max(0, len(rows) - limit)

    comment = spec.get("comment")
    width = BODY_W if not comment else BODY_W - Inches(3.5)
    avail = SLIDE_H - top - Inches(1.05)
    height = min(Inches(0.36) * (len(shown) + 1), avail)
    _add_table(slide, MARGIN, top + Inches(0.1), width, height, cols, shown,
               font=11.5 if len(cols) <= 7 else 10,
               highlight_rows=set(spec.get("highlight_rows") or []))
    if comment:
        _box(slide, MARGIN + width + Inches(0.3), top + Inches(0.1),
             Inches(3.2), avail, comment, size=13.5)
    note = spec.get("source", "")
    if cut:
        note = (note + f"（全 {len(rows):,} 行のうち上位 {limit} 行）").strip()
    _source(slide, SLIDE_H - Inches(0.85), note)
    if spec.get("callout"):
        _callout(slide, MARGIN, SLIDE_H - Inches(1.5), BODY_W, spec["callout"])
    _notes(slide, spec.get("notes", ""))
    return slide


def _style_chart(chart, spec, series_count):
    chart.has_title = False
    chart.font.size = Pt(12)
    chart.font.name = config.REPORT_FONT_JA
    kind = str(spec.get("chart") or "bar").lower()

    if series_count > 1 or kind in ("pie", "doughnut"):
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(12)
    else:
        chart.has_legend = False

    try:
        for i, s in enumerate(chart.series):
            s.format.fill.solid()
            s.format.fill.fore_color.rgb = RGBColor.from_string(
                SERIES[i % len(SERIES)])
            s.format.line.color.rgb = RGBColor.from_string(SERIES[i % len(SERIES)])
    except (AttributeError, ValueError, NotImplementedError):
        pass

    # 数値軸は桁区切り。目盛り線は薄く。
    try:
        va = chart.value_axis
        va.has_major_gridlines = True
        va.major_gridlines.format.line.color.rgb = RGBColor(0xE8, 0xE8, 0xE8)
        va.format.line.color.rgb = LINE
        va.tick_labels.number_format = spec.get("number_format") or "#,##0"
        va.tick_labels.number_format_is_linked = False
        va.tick_labels.font.size = Pt(11.5)
    except (AttributeError, ValueError, NotImplementedError):
        pass
    try:
        ca = chart.category_axis
        ca.has_major_gridlines = False
        ca.format.line.color.rgb = LINE
        ca.tick_labels.font.size = Pt(11.5)
    except (AttributeError, ValueError, NotImplementedError):
        pass

    # 値ラベル。多すぎると潰れるので、少ないときだけ。
    want = spec.get("data_labels")
    cats = len(spec.get("categories") or [])
    if want is None:
        want = kind in ("pie", "doughnut") or (cats and cats <= 8 and series_count <= 2)
    if want:
        try:
            plot = chart.plots[0]
            plot.has_data_labels = True
            dl = plot.data_labels
            dl.font.size = Pt(11)
            dl.font.name = config.REPORT_FONT_JA
            if kind in ("pie", "doughnut"):
                dl.show_percentage = True
                dl.number_format = "0.0%"
                dl.number_format_is_linked = False
                dl.position = XL_LABEL_POSITION.OUTSIDE_END
            else:
                dl.number_format = spec.get("number_format") or "#,##0"
                dl.number_format_is_linked = False
                if kind in ("bar", "hbar"):
                    dl.position = XL_LABEL_POSITION.OUTSIDE_END
        except (AttributeError, ValueError, NotImplementedError):
            pass


def _slide_chart(prs, spec):
    slide = _blank(prs)
    top = _header(slide, spec.get("title", ""), spec.get("message", ""))
    series = spec.get("series") or []
    cats = [_clean(c) for c in (spec.get("categories") or [])]

    # 画像として渡された図（PowerPointで描けない種類）はそのまま貼る
    if spec.get("image"):
        avail_h = SLIDE_H - top - Inches(1.05)
        has_side = bool(spec.get("comment"))
        w = BODY_W if not has_side else BODY_W - Inches(3.5)
        slide.shapes.add_picture(io.BytesIO(spec["image"]), MARGIN,
                                 top + Inches(0.08), width=w)
        if has_side:
            _box(slide, MARGIN + w + Inches(0.3), top + Inches(0.1), Inches(3.2),
                 avail_h, spec["comment"], size=13.5)
        _source(slide, SLIDE_H - Inches(0.85), spec.get("source", ""))
        if spec.get("callout"):
            _callout(slide, MARGIN, SLIDE_H - Inches(1.5), BODY_W, spec["callout"])
        _notes(slide, spec.get("notes", ""))
        return slide

    if not series:
        raise ReportError("chart スライドには series か image が必要です。")
    kind = str(spec.get("chart") or "bar").lower()
    if kind not in CHART_TYPES:
        raise ReportError(f"未対応のグラフ種類です: {kind}。"
                          f"使えるのは {', '.join(CHART_TYPES)} です。")
    if len(cats) > MAX_CATEGORIES:
        cats = cats[:MAX_CATEGORIES]
        series = [{**s, "values": (s.get("values") or [])[:MAX_CATEGORIES]}
                  for s in series]
        spec = {**spec, "source": (spec.get("source", "")
                                   + f"（上位{MAX_CATEGORIES}件）").strip()}

    if kind == "scatter":
        data = XyChartData()
        for s in series:
            sd = data.add_series(_clean(s.get("name") or "系列"))
            for x, y in zip(s.get("x") or [], s.get("values") or s.get("y") or []):
                if _num(x) is not None and _num(y) is not None:
                    sd.add_data_point(_num(x), _num(y))
    else:
        data = CategoryChartData()
        data.categories = cats or [str(i + 1) for i in
                                   range(len(series[0].get("values") or []))]
        for s in series:
            data.add_series(_clean(s.get("name") or "系列"),
                            [_num(v) for v in (s.get("values") or [])],
                            number_format=spec.get("number_format") or "#,##0")

    has_side = bool(spec.get("comment"))
    width = BODY_W if not has_side else BODY_W - Inches(3.5)
    height = SLIDE_H - top - Inches(1.05)
    frame = slide.shapes.add_chart(CHART_TYPES[kind], MARGIN, top + Inches(0.08),
                                   width, height, data)
    _style_chart(frame.chart, {**spec, "categories": cats}, len(series))

    if has_side:
        _box(slide, MARGIN + width + Inches(0.3), top + Inches(0.15), Inches(3.2),
             height - Inches(0.2), spec["comment"], size=13.5)
    _source(slide, SLIDE_H - Inches(0.85), spec.get("source", ""))
    if spec.get("callout"):
        _callout(slide, MARGIN, SLIDE_H - Inches(1.5), BODY_W, spec["callout"])
    _notes(slide, spec.get("notes", ""))
    return slide


def _slide_compare(prs, spec):
    """左右に並べて比べる（案A/案B、前年/今年 など）。"""
    slide = _blank(prs)
    top = _header(slide, spec.get("title", ""), spec.get("message", ""))
    panes = (spec.get("panes") or [])[:2]
    if len(panes) != 2:
        raise ReportError("compare スライドには panes を2つ指定してください。")
    gap = Inches(0.4)
    w = (BODY_W - gap) / 2
    h = SLIDE_H - top - Inches(1.05)
    for i, pane in enumerate(panes):
        left = MARGIN + (w + gap) * i
        head = _rect(slide, left, top + Inches(0.05), w, Inches(0.46),
                     fill=NAVY if i == 0 else ACCENT)
        head.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        _text(head.text_frame, _clean(pane.get("title")), size=15, bold=True,
              color=WHITE, align=PP_ALIGN.CENTER, space_after=0)
        y = top + Inches(0.62)
        if pane.get("value") is not None:
            _box(slide, left, y, w, Inches(0.8),
                 _fmt(pane["value"]) + _clean(pane.get("unit")),
                 size=30, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
            y += Inches(0.9)
        if pane.get("image"):
            slide.shapes.add_picture(io.BytesIO(pane["image"]), left, y, width=w)
            y += Inches(2.6)
        if pane.get("bullets"):
            _bullets(slide, left + Inches(0.05), y, w - Inches(0.1),
                     top + h - y, pane["bullets"], size=14)
    if spec.get("callout"):
        _callout(slide, MARGIN, SLIDE_H - Inches(1.5), BODY_W, spec["callout"])
    _source(slide, SLIDE_H - Inches(0.85), spec.get("source", ""))
    _notes(slide, spec.get("notes", ""))
    return slide


def _slide_closing(prs, spec):
    """まとめと次のアクション。担当と期限まで書けるようにする。"""
    slide = _blank(prs)
    top = _header(slide, spec.get("title") or "まとめと次のアクション",
                  spec.get("message", ""))
    y = top + Inches(0.1)
    if spec.get("summary"):
        _box(slide, MARGIN, y, BODY_W, Inches(0.3), "まとめ", size=12,
             bold=True, color=MUTED)
        y = _bullets(slide, MARGIN, y + Inches(0.34), BODY_W, Inches(2.0),
                     spec["summary"], size=16) + Inches(0.15)

    actions = spec.get("actions") or []
    if actions:
        _box(slide, MARGIN, y, BODY_W, Inches(0.3), "次のアクション", size=12,
             bold=True, color=MUTED)
        y += Inches(0.34)
        rows = []
        for a in actions:
            if isinstance(a, dict):
                rows.append([a.get("text", ""), a.get("owner", ""), a.get("due", "")])
            else:
                rows.append([str(a), "", ""])
        h = min(Inches(0.36) * (len(rows) + 1), SLIDE_H - y - Inches(0.8))
        _add_table(slide, MARGIN, y, BODY_W, h, ["やること", "担当", "期限"], rows,
                   font=12.5)
    _notes(slide, spec.get("notes", ""))
    return slide


_BUILDERS = {"title": _slide_title, "agenda": _slide_agenda,
             "section": _slide_section, "message": _slide_message,
             "table": _slide_table, "chart": _slide_chart, "kpi": _slide_kpi,
             "compare": _slide_compare, "closing": _slide_closing,
             # 旧称
             "text": _slide_message}


# =============================================================================
# 組み立て
# =============================================================================

def build(slides: list[dict], title: str | None = None,
          subtitle: str | None = None, footer: str | None = None,
          agenda: bool = True) -> bytes:
    """スライド定義のリストから .pptx のバイト列を作る。"""
    if not slides:
        raise ReportError("スライドが1枚もありません。")
    prs = Presentation()
    prs.slide_width, prs.slide_height = SLIDE_W, SLIDE_H

    specs = list(slides)
    if title and (specs[0].get("kind") or "").lower() != "title":
        specs.insert(0, {"kind": "title", "title": title, "subtitle": subtitle or ""})

    # 中扉があれば、その並びから目次を自動で作る
    if agenda and not any((s.get("kind") or "") == "agenda" for s in specs):
        sections = [s.get("title", "") for s in specs
                    if (s.get("kind") or "") == "section"]
        if len(sections) >= 2:
            at = 1 if (specs[0].get("kind") or "") == "title" else 0
            specs.insert(at, {"kind": "agenda", "title": "本日の内容",
                              "items": sections})

    page = 0
    for i, spec in enumerate(specs):
        kind = str(spec.get("kind") or "message").lower()
        if kind not in _BUILDERS:
            raise ReportError(f"{i + 1}枚目: 未対応の種類です: {kind}。"
                              f"使えるのは {', '.join(SLIDE_KINDS)} です。")
        try:
            slide = _BUILDERS[kind](prs, spec)
        except ReportError:
            raise
        except Exception as e:
            raise ReportError(f"{i + 1}枚目（{kind}）の作成に失敗しました: {e}") from e
        if kind not in ("title", "section"):
            page += 1
            _footer(slide, footer or config.REPORT_ORG or "", page)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def safe_filename(name: str | None, default: str = "report") -> str:
    base = re.sub(r'[\\/:*?"<>|]', "_", str(name or default)).strip() or default
    return base if base.lower().endswith(".pptx") else base + ".pptx"


def outline(slides: list[dict]) -> list[str]:
    """何が入ったかの一覧（画面とLLMへの報告用）。"""
    labels = {"title": "表紙", "agenda": "目次", "section": "中扉",
              "message": "説明", "text": "説明", "table": "表", "chart": "グラフ",
              "kpi": "KPI", "compare": "比較", "closing": "まとめ"}
    out = []
    for i, s in enumerate(slides, 1):
        kind = str(s.get("kind") or "message").lower()
        extra = ""
        if kind == "chart":
            extra = (f"（{s.get('chart', 'bar')}・{len(s.get('series') or [])}系列）"
                     if not s.get("image") else "（画像）")
        elif kind == "table":
            extra = f"（{len(s.get('rows') or [])}行）"
        out.append(f"{i}. [{labels.get(kind, kind)}{extra}] {s.get('title', '')}")
    return out
